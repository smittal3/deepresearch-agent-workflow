"""Deep Research Agent — Streamlit front end.

Upload a document (a pitch deck, memo, or TXT), ask a question, and four agents
(RAG / Web / Academic / Synthesis) investigate in parallel and return a
decision-ready dashboard. Powered end-to-end by a single Google Gemini key.
"""
from __future__ import annotations

import html
import io
import json
import os
import time

import streamlit as st
from dotenv import load_dotenv

from graph import AGENT_STEPS, PARALLEL_AGENTS, build_graph
from llm import MODEL, set_global_key

load_dotenv()

st.set_page_config(
    page_title="Deep Research Agent",
    page_icon="🔭",
    layout="centered",
    initial_sidebar_state="expanded",
)

# --------------------------------------------------------------------------- #
# Styling
# --------------------------------------------------------------------------- #
st.markdown(
    """
    <style>
      .hero {
        background: linear-gradient(120deg, #7C6CF0 0%, #4A3FD6 45%, #2BB6C4 100%);
        padding: 1.6rem 2rem; border-radius: 18px; margin-bottom: 1.2rem;
        box-shadow: 0 10px 30px rgba(0,0,0,.35);
      }
      .hero h1 { color: #fff; margin: 0; font-size: 2.1rem; letter-spacing: -.5px; }
      .hero p  { color: rgba(255,255,255,.88); margin: .35rem 0 0; font-size: 1.02rem; }
      .agent-row {
        display:flex; align-items:center; gap:.7rem; padding:.6rem .9rem; margin:.35rem 0;
        background:#171A23; border:1px solid #262B38; border-radius:12px;
      }
      .agent-row .name { font-weight:600; }
      .agent-row .desc { color:#9AA3B2; font-size:.86rem; }
      .badge {
        display:inline-block; padding:.28rem .8rem; border-radius:999px;
        font-weight:700; font-size:.82rem; letter-spacing:.3px;
      }
      .verdict-box {
        background:#171A23; border:1px solid #262B38; border-left:6px solid #7C6CF0;
        border-radius:14px; padding:1.1rem 1.3rem; margin:.4rem 0 1rem;
      }
      .verdict-box h3 { margin:.1rem 0 .5rem; }
      div[data-testid="stMetric"] {
        background:#171A23; border:1px solid #262B38; border-radius:14px;
        padding:.8rem 1rem;
      }
      .dp-grid {
        display:grid; grid-template-columns:repeat(auto-fit, minmax(190px, 1fr));
        gap:.7rem; margin:.2rem 0 .4rem;
      }
      .dp-card {
        background:#171A23; border:1px solid #262B38; border-radius:14px;
        padding:.85rem 1rem; display:flex; flex-direction:column;
      }
      .dp-label { color:#9AA3B2; font-size:.82rem; margin-bottom:.25rem; }
      .dp-value {
        font-size:1.5rem; font-weight:700; line-height:1.2; overflow-wrap:anywhere;
      }
      .dp-value.long { font-size:1.02rem; line-height:1.3; }
      .dp-cites { margin-top:auto; padding-top:.3rem; }
      .flag-card {
        background:#1d1614; border:1px solid #3a2420; border-left:4px solid #ef4444;
        border-radius:10px; padding:.7rem .9rem; margin:.45rem 0; color:#ecd9d5;
        font-size:.92rem; line-height:1.45;
      }
      .flag-ok {
        background:#14201a; border:1px solid #21402f; border-left:4px solid #22c55e;
        border-radius:10px; padding:.7rem .9rem; margin:.45rem 0; color:#cfe9d8;
        font-size:.92rem;
      }
      .cite {
        display:inline-block; text-decoration:none; font-size:.72rem; font-weight:700;
        background:#7C6CF022; color:#9d92f5; border:1px solid #7C6CF0; border-radius:6px;
        padding:0 .4rem; margin:.35rem .25rem 0 0;
      }
      .cite:hover { background:#7C6CF0; color:#fff; }
      .score-row { margin:.55rem 0; }
      .score-head { display:flex; justify-content:space-between; font-size:.9rem; margin-bottom:.25rem; }
      .score-track { background:#262B38; border-radius:999px; height:10px; overflow:hidden; }
      .score-fill { height:100%; border-radius:999px; }
      .score-why { color:#9AA3B2; font-size:.78rem; margin-top:.15rem; }
      .src-item { padding:.3rem 0; border-bottom:1px solid #1f2430; font-size:.88rem; }
      .src-tag {
        display:inline-block; font-size:.68rem; font-weight:700; border-radius:5px;
        padding:0 .35rem; margin-right:.45rem; text-transform:uppercase;
      }
      .step-done { color:#7f8896; font-size:.86rem; padding:.12rem 0; }
      .step-done b { color:#cdd3df; }
    </style>
    """,
    unsafe_allow_html=True,
)


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #
MIN_PDF_TEXT_CHARS = 200  # below this we assume an image-based PDF and use Gemini vision
OCR_MODEL = "gemini-2.5-flash"  # fast/cheap model for transcription


def pdf_text_via_gemini(data: bytes, api_key: str) -> str:
    """Transcribe an image-based / scanned PDF using Gemini's native PDF support.

    Pitch decks are often exported as images, so pypdf finds no text layer. Since we
    already depend on Gemini, we let it read the PDF directly — no OCR system deps."""
    try:
        from google import genai
        from google.genai import types

        client = genai.Client(api_key=api_key)
        resp = client.models.generate_content(
            model=OCR_MODEL,
            contents=[
                types.Part.from_bytes(data=data, mime_type="application/pdf"),
                "Transcribe ALL text and information from this document verbatim — "
                "including text on slides, in charts, tables, and figures. Preserve "
                "numbers and labels exactly. Output plain text only, no commentary.",
            ],
        )
        return (resp.text or "").strip()
    except Exception:  # noqa: BLE001
        return ""


def extract_text(uploaded_file, api_key: str = "") -> str:
    """Pull plain text out of an uploaded file (Streamlit UploadedFile)."""
    if uploaded_file is None:
        return ""
    return extract_text_from_bytes(uploaded_file.name, uploaded_file.getvalue(), api_key)


def extract_text_from_path(path: str, api_key: str = "") -> str:
    """Pull plain text out of a file on disk (used for bundled samples)."""
    with open(path, "rb") as f:
        return extract_text_from_bytes(os.path.basename(path), f.read(), api_key)


def extract_text_from_bytes(filename: str, data: bytes, api_key: str = "") -> str:
    """Core extractor for PDF, DOCX, PPTX, or TXT/MD content."""
    name = filename.lower()

    if name.endswith(".pdf"):
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        text = "\n".join(page.extract_text() or "" for page in reader.pages).strip()
        if len(text) < MIN_PDF_TEXT_CHARS and api_key:
            ocr = pdf_text_via_gemini(data, api_key)
            if len(ocr) > len(text):
                return ocr
        return text

    if name.endswith(".docx"):
        from docx import Document as DocxDocument

        doc = DocxDocument(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:  # capture tabular financials too
            for row in table.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(t for t in parts if t.strip()).strip()

    if name.endswith(".pptx"):
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [
                shape.text for shape in slide.shapes
                if shape.has_text_frame and shape.text.strip()
            ]
            if texts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides).strip()

    # txt / md / anything text-like
    return data.decode("utf-8", errors="ignore").strip()


def report_to_markdown(query: str, r: dict) -> str:
    lines = [
        "# Deep Research Report",
        f"\n**Query:** {query}\n",
        f"**Verdict:** {r.get('bottom_line', '')}  ",
        f"**Confidence:** {r.get('confidence', '')}\n",
        "## Executive Summary",
        r.get("executive_summary", ""),
        "\n## Key Data Points",
    ]
    for dp in r.get("key_data_points", []):
        cites = "".join(f" [{i}]" for i in dp.get("source_ids", []))
        lines.append(f"- **{dp['label']}:** {dp['value']}{cites}")

    scorecard = r.get("scorecard", [])
    if scorecard:
        lines += ["\n## Scorecard"]
        for item in scorecard:
            lines.append(
                f"- **{item.get('dimension','')}: {item.get('score','')}/5** — "
                f"{item.get('rationale','')}"
            )

    lines += ["\n## Market Sentiment", r.get("sentiment_analysis", "")]
    lines += ["\n## Technical Reality", r.get("technical_reality", "")]
    lines += ["\n## Contradictions Found"]
    contradictions = r.get("contradictions_found", [])
    if contradictions:
        lines += [f"- ⚠️ {c}" for c in contradictions]
    else:
        lines.append("- None detected — internal claims align with external evidence.")
    lines += ["\n## Recommended Due-Diligence Questions"]
    lines += [f"- {q}" for q in r.get("recommended_questions", [])]

    sources = r.get("sources", [])
    if sources:
        lines += ["\n## Sources"]
        for s in sources:
            url = f" — {s['url']}" if s.get("url") else ""
            lines.append(f"- [{s['id']}] ({s.get('type','')}) {s.get('title','')}{url}")
    return "\n".join(lines)


SAMPLES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "samples")
SAMPLE_EXTS = (".pdf", ".docx", ".pptx", ".txt", ".md")


def _list_files_in(dir_path: str) -> list[str]:
    """Loadable sample filenames in a directory (README/manifest excluded)."""
    if not os.path.isdir(dir_path):
        return []
    return [
        f for f in sorted(os.listdir(dir_path))
        if f.lower().endswith(SAMPLE_EXTS) and f.lower() != "readme.md"
    ]


# Used if samples.json is missing/unreadable so the app still offers one case.
_DEFAULT_CASE = {
    "id": "startup",
    "label": "Startup due-diligence",
    "emoji": "🚀",
    "dir": "startup_due_diligence",
    "blurb": "VC evaluating a seed-stage startup.",
    "query": (
        "Act as a Venture Capitalist. I've uploaded the investment and technical files "
        "for a startup we are evaluating. Please run a deep due-diligence report on their "
        "core technology and market claims. Verify whether their product claims are "
        "scientifically valid, and give us a recommendation on whether we should invest "
        "in this startup at this $5M seed stage."
    ),
}


def list_sample_cases() -> list[dict]:
    """Selectable sample cases from samples/samples.json.

    Each case = {id, label, emoji, dir, blurb, query, files}; `files` is filled from the
    case's subdirectory. Cases with no readable files are dropped. Falls back to a single
    case (or a legacy flat `case_query` manifest) so the app never loses the sample button.
    """
    manifest = os.path.join(SAMPLES_DIR, "samples.json")
    cases: list[dict] = []
    if os.path.exists(manifest):
        try:
            with open(manifest) as f:
                data = json.load(f)
            if isinstance(data, dict) and data.get("cases"):
                cases = list(data["cases"])
            elif isinstance(data, dict) and data.get("case_query"):
                # Legacy single-case manifest: treat flat samples/ as one case.
                cases = [{**_DEFAULT_CASE, "dir": "", "query": data["case_query"]}]
        except Exception:  # noqa: BLE001
            cases = []
    if not cases:
        cases = [_DEFAULT_CASE]

    out: list[dict] = []
    for c in cases:
        case = dict(c)
        case["files"] = _list_files_in(os.path.join(SAMPLES_DIR, case.get("dir", "")))
        if case["files"]:
            out.append(case)
    return out


def load_sample_documents(case: dict, api_key: str = "") -> tuple[str, list[str]]:
    """Read every file in a sample case's directory into one combined blob + names."""
    base = os.path.join(SAMPLES_DIR, case.get("dir", ""))
    texts, names = [], []
    for fname in case.get("files") or _list_files_in(base):
        try:
            text = extract_text_from_path(os.path.join(base, fname), api_key)
        except Exception:  # noqa: BLE001
            text = ""
        if text.strip():
            texts.append(f"===== FILE: {fname} =====\n{text}")
            names.append(fname)
    return "\n\n".join(texts), names


CONF_COLORS = {"high": "#22c55e", "medium": "#f59e0b", "low": "#ef4444"}


# --------------------------------------------------------------------------- #
# Sidebar — key + model
# --------------------------------------------------------------------------- #
with st.sidebar:
    st.header("⚙️ Setup")

    env_key = os.getenv("GOOGLE_API_KEY", "")
    if env_key:
        # Key already provided via .env / environment — no need to paste anything.
        st.success("✓ API key detected from `.env`")
        override = st.text_input(
            "Use a different key (optional)",
            type="password",
            placeholder="leave blank to use .env key",
        )
        api_key = override.strip() or env_key
    else:
        api_key = st.text_input(
            "Google Gemini API key",
            type="password",
            help="One key powers all four agents + embeddings.",
            placeholder="AIza...",
        ).strip()
        with st.expander("🔑 How to get a free key"):
            st.markdown(
                "1. Go to **[aistudio.google.com/app/apikey]"
                "(https://aistudio.google.com/app/apikey)**\n"
                "2. Click **Create API key**\n"
                "3. Paste it above, or add it to a `.env` file as `GOOGLE_API_KEY=...`"
            )

    st.caption(f"Model: **{MODEL}**")
    st.caption("Your key stays in this session only — it is never stored.")
    st.divider()
    st.markdown("**Pipeline**")
    for _, name, desc in AGENT_STEPS:
        st.markdown(f"{name}  \n<span style='color:#9AA3B2;font-size:.8rem'>{desc}</span>",
                    unsafe_allow_html=True)


# --------------------------------------------------------------------------- #
# Phase machine — one screen at a time: input -> running -> results
# --------------------------------------------------------------------------- #
if "query_text" not in st.session_state:
    st.session_state.query_text = ""
phase = st.session_state.get("phase", "input")


def goto(new_phase: str):
    st.session_state.phase = new_phase
    st.rerun()


HERO = """
    <div class="hero">
      <h1>🔭 Deep Research Agent</h1>
      <p>Four agents read your documents, search the web, check academic papers, and
      cross-examine the findings into a cited report — working in parallel.</p>
    </div>
"""

# ---------------------------- INPUT SCREEN --------------------------------- #
if phase == "input":
    st.markdown(HERO, unsafe_allow_html=True)

    sample_cases = list_sample_cases()
    if sample_cases:
        st.caption("Try a bundled case study, or upload your own documents below.")
        cols = st.columns(len(sample_cases))
        for col, case in zip(cols, sample_cases):
            label = f"{case.get('emoji', '📂')}  {case.get('label', 'Sample case')}"
            if col.button(label, use_container_width=True, help=case.get("blurb", "")):
                st.session_state.sample_loaded = True
                st.session_state.sample_case_id = case["id"]
                st.session_state.query_text = case.get("query", "")
                st.rerun()

    st.subheader("Documents (optional)")
    uploaded = st.file_uploader(
        "PDF, DOCX, PPTX, TXT — add several if you like",
        type=["pdf", "docx", "pptx", "txt", "md"],
        accept_multiple_files=True,
        label_visibility="collapsed",
    )
    if uploaded:
        st.session_state.sample_loaded = False  # uploads take precedence over the sample
        st.caption("📎 Using " + ", ".join(f"**{u.name}**" for u in uploaded))
    elif st.session_state.get("sample_loaded"):
        loaded = next(
            (c for c in sample_cases if c["id"] == st.session_state.get("sample_case_id")),
            None,
        )
        if loaded:
            cc = st.columns([4, 1])
            cc[0].caption(
                f"📎 {loaded.get('emoji', '')} {loaded['label']}: "
                + ", ".join(loaded["files"])
            )
            if cc[1].button("✕ Clear", use_container_width=True):
                st.session_state.sample_loaded = False
                st.rerun()
        else:
            st.session_state.sample_loaded = False

    st.subheader("Research query")
    query = st.text_area(
        "Research query", key="query_text", height=150, label_visibility="collapsed",
        placeholder="Describe what you want to learn, compare, or evaluate...",
    )
    st.caption("Tip: name the key companies, technologies, or concepts so the agents "
               "know what to search for.")

    if st.button("🚀  Run Deep Research", type="primary", use_container_width=True):
        if not api_key:
            st.error("⚠️ Add your Google Gemini API key in the sidebar to run.")
        elif not query.strip():
            st.error("⚠️ Enter a research query first.")
        else:
            # Stash inputs — the widgets won't exist on the running screen.
            st.session_state.pending_uploads = (
                [(u.name, u.getvalue()) for u in uploaded] if uploaded else []
            )
            st.session_state.pending_sample = bool(
                st.session_state.get("sample_loaded") and not uploaded
            )
            st.session_state.pending_sample_case_id = st.session_state.get("sample_case_id")
            st.session_state.pending_query = query.strip()
            st.session_state.pop("report", None)
            goto("running")


# --------------------------------------------------------------------------- #
# Report renderer (single clean page; details live in expanders)
# --------------------------------------------------------------------------- #
def render_report(query: str, r: dict):
    conf = (r.get("confidence", "Medium") or "Medium").strip()
    color = CONF_COLORS.get(conf.lower(), "#7C6CF0")
    sources = r.get("sources", [])
    src_by_id = {s["id"]: s for s in sources}
    contradictions = r.get("contradictions_found", [])

    def cite_chips(ids):
        chips = []
        for i in ids:
            s = src_by_id.get(i)
            if not s:
                continue
            if s.get("url"):
                chips.append(f"<a class='cite' href='{s['url']}' target='_blank'>[{i}]</a>")
            else:
                chips.append(f"<span class='cite'>[{i}]</span>")  # document, no link
        return "".join(chips)

    # ---- Verdict (prominent) ----
    st.markdown(
        f"""<div class="verdict-box">
          <span class="badge" style="background:{color}22;color:{color};border:1px solid {color}">
            {conf.upper()} CONFIDENCE</span>
          <h3 style="margin-top:.6rem">{r.get('bottom_line','')}</h3>
        </div>""",
        unsafe_allow_html=True,
    )
    st.markdown(r.get("executive_summary", ""))

    # ---- Scorecard ----
    scorecard = r.get("scorecard", [])
    if scorecard:
        st.markdown("##### 🎯 Scorecard")
        for item in scorecard:
            score = max(1, min(5, int(item.get("score", 0) or 0)))
            sc_color = ["#ef4444", "#f97316", "#f59e0b", "#84cc16", "#22c55e"][score - 1]
            st.markdown(
                f"<div class='score-row'>"
                f"<div class='score-head'><span>{item.get('dimension','')}</span>"
                f"<span style='color:{sc_color};font-weight:700'>{score}/5</span></div>"
                f"<div class='score-track'><div class='score-fill' "
                f"style='width:{score*20}%;background:{sc_color}'></div></div>"
                f"<div class='score-why'>{item.get('rationale','')}</div></div>",
                unsafe_allow_html=True,
            )

    # ---- Key data points ----
    data_points = r.get("key_data_points", [])
    if data_points:
        st.markdown("##### 📊 Key data points")
        cards = []
        for dp in data_points:
            value = str(dp.get("value", ""))
            long = " long" if len(value) > 22 else ""  # shrink the font for sentence-y values
            chips = cite_chips(dp.get("source_ids", []))
            cards.append(
                f"<div class='dp-card'>"
                f"<div class='dp-label'>{dp.get('label','')}</div>"
                f"<div class='dp-value{long}'>{value}</div>"
                f"<div class='dp-cites'>{chips}</div></div>"
            )
        st.markdown(f"<div class='dp-grid'>{''.join(cards)}</div>", unsafe_allow_html=True)

    # ---- Red flags (the value-add, kept visible) ----
    st.markdown("##### ⚠️ Contradictions & red flags")
    if contradictions:
        st.markdown(
            "".join(
                f"<div class='flag-card'>⚠️&nbsp; {html.escape(str(c))}</div>"
                for c in contradictions
            ),
            unsafe_allow_html=True,
        )
    else:
        st.markdown(
            "<div class='flag-ok'>✓&nbsp; No contradictions detected — the sources are "
            "consistent with each other.</div>",
            unsafe_allow_html=True,
        )

    # ---- Detailed analysis (open by default, but collapsible) ----
    with st.expander("📝 Detailed analysis", expanded=True):
        st.markdown("**🌐 Web & market sentiment**")
        st.markdown(r.get("sentiment_analysis", "_(none)_"))
        st.markdown("**🔬 Technical & academic reality**")
        st.markdown(r.get("technical_reality", "_(none)_"))
        questions = r.get("recommended_questions", [])
        if questions:
            st.markdown("**❓ Questions to investigate next**")
            for q in questions:
                st.markdown(f"- {q}")

    with st.expander(f"🔗 Sources ({len(sources)})"):
        if sources:
            tag_colors = {"web": "#2BB6C4", "academic": "#a78bfa", "document": "#f59e0b"}
            rows = []
            for s in sources:
                c = tag_colors.get(s.get("type", ""), "#7C6CF0")
                tag = (
                    f"<span class='src-tag' style='background:{c}22;color:{c};"
                    f"border:1px solid {c}'>{s.get('type','')}</span>"
                )
                title = s.get("title", "") or s.get("url", "")
                body = (
                    f"<a href='{s['url']}' target='_blank' style='color:#cdd3df'>{title}</a>"
                    if s.get("url") else title
                )
                rows.append(f"<div class='src-item'><b>[{s['id']}]</b> {tag}{body}</div>")
            st.markdown("".join(rows), unsafe_allow_html=True)
        else:
            st.caption("No external sources were retrieved for this run.")

    notes = r.get("_agent_notes") or {}
    if any(notes.values()):
        with st.expander("🔍 What each agent found"):
            for key, label in [("rag", "📄 Document Analyst"),
                               ("web", "🌐 Web Researcher"),
                               ("academic", "🔬 Academic Validator")]:
                if notes.get(key):
                    st.markdown(f"**{label}**")
                    st.markdown(notes[key])
    with st.expander("🔧 Raw JSON"):
        st.json({k: v for k, v in r.items() if k != "_agent_notes"})


# --------------------------------- RUNNING SCREEN -------------------------- #
if phase == "running":
    st.markdown(
        "<h3 style='text-align:center;margin-top:1.5rem'>🔭 Researching…</h3>"
        "<p style='text-align:center;color:#9AA3B2;margin-top:-.4rem'>Four agents are "
        "working in parallel — this usually takes about a minute.</p>",
        unsafe_allow_html=True,
    )
    set_global_key(api_key)
    start = time.time()
    final_report, collected = None, {}
    try:
        with st.status("Reading your documents…", expanded=True) as status:
            # Rebuild documents from the stashed inputs.
            document_text, names = "", []
            uploads = st.session_state.get("pending_uploads", [])
            if uploads:
                texts = []
                for fname, data in uploads:
                    t = extract_text_from_bytes(fname, data, api_key)
                    if t.strip():
                        texts.append(f"===== FILE: {fname} =====\n{t}")
                        names.append(fname)
                document_text = "\n\n".join(texts)
            elif st.session_state.get("pending_sample"):
                case = next(
                    (c for c in list_sample_cases()
                     if c["id"] == st.session_state.get("pending_sample_case_id")),
                    None,
                )
                if case:
                    document_text, names = load_sample_documents(case, api_key)
            if names:
                status.write(f"✓ Read {len(names)} document(s): {', '.join(names)}")

            initial_state = {
                "user_query": st.session_state.get("pending_query", ""),
                "document_text": document_text,
                "has_document": bool(document_text),
                "doc_name": ", ".join(names),
                "doc_names": names,
                "api_key": api_key,
                "model": MODEL,
            }

            status.update(label="🛰️ Agents researching in parallel…")
            done_parallel = 0
            for chunk in build_graph().stream(initial_state):
                for node, update in chunk.items():
                    update = update or {}
                    collected.update(update)
                    if node == "rag_agent" and initial_state["has_document"]:
                        status.write(f"✓ Documents analyzed — {len(names)} file(s)")
                    elif node == "web_agent":
                        status.write(f"✓ Web research — {len(update.get('web_sources', []))} sources")
                    elif node == "academic_agent":
                        status.write(f"✓ Academic search — {len(update.get('academic_sources', []))} papers")
                    elif node == "synthesis_agent":
                        final_report = update.get("final_report")
                        status.write("✓ Report synthesized")
                    if node in PARALLEL_AGENTS:
                        done_parallel += 1
                        if done_parallel == len(PARALLEL_AGENTS):
                            status.update(label="🧠 Cross-examining sources & writing the report…")
            status.update(label="✓ Research complete", state="complete")
    except Exception as e:  # noqa: BLE001
        st.error(f"Research run failed: {e}")
        if st.button("← Back"):
            goto("input")
        st.stop()

    if not final_report:
        st.error("No report was produced. Check your API key and try again.")
        if st.button("← Back"):
            goto("input")
        st.stop()

    final_report["_agent_notes"] = {
        "rag": collected.get("rag_notes", ""),
        "web": collected.get("web_notes", ""),
        "academic": collected.get("academic_notes", ""),
    }
    st.session_state.report = final_report
    st.session_state.report_query = st.session_state.get("pending_query", "")
    st.session_state.report_elapsed = time.time() - start
    goto("results")


# --------------------------------- RESULTS SCREEN -------------------------- #
if phase == "results":
    report = st.session_state.get("report")
    if not report:
        goto("input")
    bar = st.columns([3, 1])
    bar[0].caption(f"✅ Completed in {st.session_state.get('report_elapsed', 0):.0f}s")
    with bar[1]:
        st.download_button(
            "⬇️ .md", data=report_to_markdown(st.session_state.report_query, report),
            file_name="deep_research_report.md", mime="text/markdown",
            use_container_width=True,
        )
    if st.button("↻  New research", use_container_width=True):
        goto("input")
    st.divider()
    render_report(st.session_state.report_query, report)
